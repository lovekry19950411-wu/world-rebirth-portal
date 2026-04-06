# -*- coding: utf-8 -*-
"""由 haris-ppt-outline.md 產生 haris-deck.pptx（可選 assets/deck 圖片）。"""
from __future__ import annotations

import re
import sys
from datetime import datetime
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
MD_PATH = ROOT / "haris-ppt-outline.md"
OUT_PATH = ROOT / "haris-deck.pptx"
FONT_NAME = "Microsoft JhengHei"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def clean_md(s: str) -> str:
    s = re.sub(r"\*\*(.+?)\*\*", r"\1", s)
    s = re.sub(r"\*(.+?)\*", r"\1", s)
    return s.strip()


def display_title(raw: str) -> str:
    if "｜" in raw:
        return raw.split("｜", 1)[-1].strip()
    return raw.strip()


def set_run_font(paragraph, size_pt: int = 18):
    for run in paragraph.runs:
        run.font.name = FONT_NAME
        run.font.size = Pt(size_pt)


def apply_font_to_shape(shape, size_pt: int = 18):
    if not shape.has_text_frame:
        return
    for p in shape.text_frame.paragraphs:
        for run in p.runs:
            run.font.name = FONT_NAME
            run.font.size = Pt(size_pt)


def parse_slides(text: str) -> list[tuple[str, list[str], list[list[str]], Path | None]]:
    chunks = re.split(r"^## ", text, flags=re.MULTILINE)
    out: list[tuple[str, list[str], list[list[str]], Path | None]] = []
    for chunk in chunks[1:]:
        lines = chunk.splitlines()
        title = lines[0].strip() if lines else ""
        img_path: Path | None = None
        bullets: list[str] = []
        table: list[list[str]] = []
        in_table = False
        for line in lines[1:]:
            raw = line.rstrip()
            mimg = re.match(r"^\[img\](.+)$", raw.strip())
            if mimg:
                rel = mimg.group(1).strip()
                img_path = (ROOT / rel).resolve()
                continue
            if raw.strip() == "---":
                continue
            if raw.startswith("|"):
                cells = [clean_md(c.strip()) for c in raw.split("|")[1:-1]]
                if len(cells) >= 2 and re.match(r"^[-:\s]+$", cells[0] or ""):
                    continue
                if cells:
                    table.append(cells)
                    in_table = True
                continue
            if in_table and not raw.startswith("|"):
                in_table = False
            if raw.startswith("- "):
                bullets.append(clean_md(raw[2:]))
            elif raw.startswith("> "):
                bullets.append(clean_md(raw[2:]))
            elif raw.strip() and not raw.startswith("#"):
                if not in_table:
                    bullets.append(clean_md(raw))
        out.append((title, bullets, table, img_path))
    return out


def resolve_image(p: Path | None) -> Path | None:
    if not p:
        return None
    if p.exists():
        return p
    for ext in (".png", ".jpg", ".jpeg", ".webp"):
        alt = p.with_suffix(ext)
        if alt.exists():
            return alt
    return None


def add_cover(prs: Presentation, main: str, rest: list[str], img: Path | None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    img = resolve_image(img)
    if img:
        with Image.open(img) as im:
            iw, ih = im.size
        # 上方主視覺，維持比例置中
        box_w = float(SLIDE_W)
        box_h = float(Inches(5.0))
        scale = min(box_w / iw, box_h / ih)
        pic_w, pic_h = iw * scale, ih * scale
        left = (float(SLIDE_W) - pic_w) / 2
        top_img = float(Inches(0.35))
        slide.shapes.add_picture(str(img), left, top_img, width=pic_w, height=pic_h)
    tx = slide.shapes.add_textbox(Inches(0.6), Inches(5.45), Inches(12.2), Inches(1.9))
    tf = tx.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = main
    set_run_font(p0, 32)
    for line in rest[:4]:
        p = tf.add_paragraph()
        p.text = line
        set_run_font(p, 18)
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.name = FONT_NAME


def add_content_slide(
    prs: Presentation,
    title: str,
    bullets: list[str],
    img: Path | None,
):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tit = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.3), Inches(0.9))
    tit.text_frame.text = title
    apply_font_to_shape(tit, 26)

    img = resolve_image(img)
    left_margin = Inches(0.55)
    text_top = Inches(1.15)
    text_w = Inches(7.0) if img else Inches(12.0)

    if img:
        with Image.open(img) as im:
            iw, ih = im.size
        pic_max_w = float(Inches(5.0))
        pic_max_h = float(Inches(5.5))
        scale = min(pic_max_w / iw, pic_max_h / ih)
        pic_w, pic_h = iw * scale, ih * scale
        pic_left = float(SLIDE_W) - pic_w - float(Inches(0.5))
        pic_top = float(Inches(1.2))
        slide.shapes.add_picture(str(img), pic_left, pic_top, width=pic_w, height=pic_h)

    body = slide.shapes.add_textbox(left_margin, text_top, text_w, Inches(5.8))
    tf = body.text_frame
    tf.word_wrap = True
    tf.clear()
    use = bullets if bullets else [" "]
    first_b = True
    for b in use:
        p = tf.paragraphs[0] if first_b else tf.add_paragraph()
        first_b = False
        p.text = b
        p.level = 0
        set_run_font(p, 17)


def add_table_slide(prs: Presentation, title: str, bullets: list[str], rows: list[list[str]]):
    if not rows or len(rows) < 2:
        add_content_slide(prs, title, bullets or [" "], None)
        return
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tit = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.85))
    tit.text_frame.text = title
    apply_font_to_shape(tit, 24)

    nrows = len(rows)
    ncols = max(len(r) for r in rows)
    tbl_left = Inches(0.45)
    tbl_top = Inches(1.05)
    tbl_w = Inches(12.3)
    table_shape = slide.shapes.add_table(nrows, ncols, tbl_left, tbl_top, tbl_w, Inches(0.42 * nrows))
    tbl = table_shape.table
    for ri, row in enumerate(rows):
        for ci in range(ncols):
            cell = tbl.cell(ri, ci)
            val = row[ci] if ci < len(row) else ""
            cell.text = val
            for p in cell.text_frame.paragraphs:
                set_run_font(p, 11)

    if bullets:
        tbox = slide.shapes.add_textbox(Inches(0.5), Inches(5.85), Inches(12.0), Inches(1.35))
        tf = tbox.text_frame
        tf.word_wrap = True
        first_b = True
        for b in bullets:
            p = tf.paragraphs[0] if first_b else tf.add_paragraph()
            first_b = False
            p.text = b
            set_run_font(p, 14)


def main() -> int:
    if not MD_PATH.exists():
        print(f"找不到: {MD_PATH}", file=sys.stderr)
        return 1
    text = MD_PATH.read_text(encoding="utf-8")
    slides_data = parse_slides(text)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    first = True
    for title, bullets, table, img in slides_data:
        t = display_title(title)
        if first and "封面" in title:
            main = bullets[0] if bullets else "華瑞絲俱樂部 HARIS"
            rest = bullets[1:] if len(bullets) > 1 else []
            add_cover(prs, main, rest, img)
            first = False
            continue
        first = False
        if table and len(table) >= 2:
            add_table_slide(prs, t, bullets, table)
        else:
            add_content_slide(prs, t, bullets, img)

    out = OUT_PATH
    try:
        prs.save(out)
    except PermissionError:
        alt = ROOT / f"haris-deck-{datetime.now().strftime('%Y%m%d-%H%M%S')}.pptx"
        prs.save(alt)
        out = alt
        print(
            "注意：haris-deck.pptx 正被其他程式開啟或鎖定，已另存新檔。請先關閉 WPS／PowerPoint 再重跑可覆寫原檔。",
            file=sys.stderr,
        )
    print(f"已產生: {out}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
